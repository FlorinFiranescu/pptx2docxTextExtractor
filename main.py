import os

import pdfplumber
import pdf2image
import docx
import pptx
from docx.shared import Pt
from docx.shared import Inches, Cm
from PIL import Image
from pptx import Presentation
import tkinter as tk
from tkinter import filedialog
from pathlib import Path

def valid_xml_char_ordinal(c):
    '''

    :param c: one-string character to convert
    :return: true if the character's unicode is valid_xml, false if not
    '''
    codepoint = ord(c)
    # conditions ordered by presumed frequency
    return (
        0x20 <= codepoint <= 0xD7FF or
        codepoint in (0x9, 0xA, 0xD) or
        0xE000 <= codepoint <= 0xFFFD or
        0x10000 <= codepoint <= 0x10FFFF
        )

def getShapesWithText(slide, maxTitleWords=10):
    '''
    Get the number of shapes with text.
    The first shape is considered to be the title
    If the title has multiple lines, we consider it a text with shape
    If the text of the other shape is not equal to an empty string consider it a shape with text

    :param slide: pptx.Presentation(pptx).slide object
    :param maxTitleWords: variable to specify the max count of words of the title.
                          If maxTitleWords >= number of words from inside the title, add it to the number of shapes
    :return: the number of the shapes with text from the slide

    a single line but it mhas ore words than specified, consider it a shape with text
    defaulted to 10
    '''
    it = 0
    for shape in slide.shapes:
        if(hasattr(shape, "text")):
            if(shape.text.strip() != ''):
                if (shape == slide.shapes[0] and len(shape.text.split('\n')) <= 1 and len(shape.text) <= maxTitleWords):
                    continue
                it += 1
    return it

def paragraphStyleSwitcher(docParagraph, pptParagraphLevel):
    '''
    use switcher to get the values
    for the style and fontSize
    work directly on the paragraph object
    no object returned

    :param docParagraph: reference to docx.Doc().paragraph object, style is set to it
    :param pptParagraphLevel:  reference to a paragraph.level integer, in rage of (0,3)
    :return: void function, works directly on the docParagraph obj
    '''
    styleSwitcher = {
        0: 'Normal',
        1: 'List Bullet',
        2: 'List Bullet 2',
        3: 'List Bullet 3'
    }
    fontSizeSwitcher = {
        0: Pt(16),
        1: Pt(15),
        2: Pt(14),
        3: Pt(13)
    }
    docParagraph.style = styleSwitcher.get(pptParagraphLevel, 'Normal')
    font = docParagraph.style.font
    font.size = fontSizeSwitcher.get(pptParagraphLevel, Pt(13))


root = tk.Tk()
root.withdraw()

pptToScan = filedialog.askopenfilename(
    title='Choose the pptx file you want to convert',
    initialdir='/',
    filetypes=[('Ppt files', '*.pptx')]
)
if not(pptToScan):
    print("No pptx file has been selected, exiting..")
    exit(0)

nameOfDoc = Path(pptToScan).stem
parentDirectory = Path(pptToScan).parent



doc = docx.Document()
style = doc.styles['List Bullet']
font = style.font
font.size = Pt(14)
docxToSave = f"{parentDirectory}/{nameOfDoc}.docx"
print(docxToSave)
print(pptToScan)
fileType = "pptx"

prs = Presentation(pptToScan)
text = ""
it = 0
for slide in prs.slides:
    it = it + 1
    print(f"Processing slide number {it} out of {len(prs.slides)}")
    """
    We want to be able to avoid slides where only
    The title and/or images are appering
    """
    if(getShapesWithText(slide) == 0):
         continue
    for shape in slide.shapes:
         if hasattr(shape, "text"):
            """
            We assume that the first shape is the title
            only if it consists of a single line. If multiple
            lines exists -> shape is supposed to be a paragraph
            """
            if (shape == slide.shapes[0] and len(shape.text.split('\n')) == 1):
               doc.add_heading(''.join(c for c in shape.text if valid_xml_char_ordinal(c)))
               continue
            for paragraph in shape.text_frame.paragraphs:
                if (paragraph.text.strip() == ''):
                    continue
                p = doc.add_paragraph()
                '''
                We work directly with the objects
                no return function required
                 '''
                paragraphStyleSwitcher(p, paragraph.level)
                runText = ""
                for run in paragraph.runs:
                    runner = p.add_run(''.join(c for c in run.text if valid_xml_char_ordinal(c)))
                    if(run.font.bold):
                        runner.bold = True
                    if(run.font.italic):
                        runner.italic = True
doc.save(docxToSave)
if(Path(docxToSave).exists()):
    print(f"Finished processing. Word document path: {docxToSave}")
    if(input("Do you wish to open the saved word? [ Y/N ]\n").lower() == "y"):
        os.system(f'start {docxToSave}')
else:
    print(f"Problem appeared while converting the ppt to word. Contact admin")

